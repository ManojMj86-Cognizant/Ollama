
# THIS SCRIPT IS NOT WORKING VERY WELL, NEED REVISION

from pptx import Presentation

def extract_text_from_pptx_robust(pptx_path):
    """
    Extracts all text from a PPTX file, with slide separation.
    """
    try:
        presentation = Presentation(pptx_path)
        all_text = ""

        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text
                        slide_text += "\n"  # Paragraph separation
            all_text += slide_text + "\n--- Slide Break ---\n"  # Slide separation

        return all_text

    except Exception as e:
        print(f"Error extracting text: {e}")
        return None

# Example usage:
if __name__ == "__main__":
    pptx_file = "./documents/Centrica Case Study.pptx"
    extracted_text = extract_text_from_pptx_robust(pptx_file)
    if extracted_text:
        print(extracted_text)