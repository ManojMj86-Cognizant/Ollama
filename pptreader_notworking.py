
# # THIS SCRIPT IS NOT WORKING VERY WELL, NEED REVISION

# from pptx import Presentation # pip install python-pptx

# def extract_text_from_pptx_robust(pptx_path):
#     """
#     Extracts all text from a PPTX file, with slide separation.
#     """
#     try:
#         presentation = Presentation(pptx_path)
#         all_text = ""

#         for slide in presentation.slides:
#             slide_text = ""
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     for paragraph in shape.text_frame.paragraphs:
#                         for run in paragraph.runs:
#                             slide_text += run.text
#                         slide_text += "\n"  # Paragraph separation
#             all_text += slide_text + "\n--- Slide Break ---\n"  # Slide separation

#         return all_text

#     except Exception as e:
#         print(f"Error extracting text: {e}")
#         return None

# # Example usage:
# if __name__ == "__main__":
#     pptx_file = "./documents/Centrica Case Study.pptx"
#     extracted_text = extract_text_from_pptx_robust(pptx_file)
#     if extracted_text:
#         print(extracted_text)

from pptx import Presentation
import os

def extract_slides_to_separate_files(pptx_path, output_dir="slide_texts"):
    """
    Extracts text from each slide of a PPTX file and saves it to separate text files.
    """
    try:
        presentation = Presentation(pptx_path)

        # Create output directory if it doesn't exist
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        for i, slide in enumerate(presentation.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text
                        slide_text += "\n"  # Paragraph separation

            # Save slide text to a separate file
            slide_filename = os.path.join(output_dir, f"slide_{i + 1}.txt")
            with open(slide_filename, "w", encoding="utf-8") as text_file:
                text_file.write(slide_text)

        return f"Text from {len(presentation.slides)} slides extracted and saved to '{output_dir}'"

    except Exception as e:
        return f"Error extracting slides: {e}"

# Example usage:
if __name__ == "__main__":
    pptx_file = "./documents/Centrica and British Gas – Account Info and Case Studies.pptx"  # Replace with your PPTX file path
    output_directory = "extracted_slides" #Name of output directory

    result = extract_slides_to_separate_files(pptx_file, output_directory)
    print(result)